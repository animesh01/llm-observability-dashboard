"""Report builders for the observability dashboard.

Produces three artifacts from the telemetry summary:
  - exec_summary_text(): plain-text executive summary (used for copy-to-clipboard / email)
  - build_pdf(): a one-page PDF report (reportlab)
  - build_pptx(): a short slide deck, one slide per bucket (python-pptx)

All inputs are the already-computed summary dict, so these functions do no
telemetry math themselves.
"""
from __future__ import annotations

import io
from datetime import date


# ----------------------------------------------------------------- text
def exec_summary_text(s: dict) -> str:
    """Plain-text executive summary suitable for email/clipboard."""
    L = []
    L.append("CONVERSATIONAL AI — MODEL HEALTH: EXECUTIVE SUMMARY")
    L.append(f"Reporting window: last {s['window_days']} days  |  Generated: {date.today().isoformat()}")
    L.append("")
    L.append(f"OVERALL STATUS: {s['status'].upper()}")
    L.append(s['headline'])
    L.append("")
    L.append("KEY NUMBERS (latest day)")
    L.append(f"  - Conversations:        {s['volume']:,}")
    L.append(f"  - Intent accuracy:      {s['intent_accuracy']*100:.1f}%")
    L.append(f"  - Resolution rate:      {s['resolution_rate']*100:.0f}%")
    L.append(f"  - Hallucination risk:   {s['hallucination_rate']*100:.1f}%")
    L.append(f"  - Groundedness:         {s['groundedness_rate']*100:.0f}%")
    L.append(f"  - p95 latency:          {s['p95_latency_ms']:,} ms")
    L.append(f"  - Cost / resolved:      ${s['cost_per_resolved']:.4f}")
    L.append(f"  - Intent drift:         {s['intent_drift']*100:.0f}%")
    L.append("")
    L.append("BY BUCKET")
    for b in s["buckets"]:
        L.append(f"  {b['name']}: {b['verdict']}")
    L.append("")
    L.append("WHAT CHANGED")
    for c in s["changes"]:
        L.append(f"  - {c}")
    L.append("")
    L.append("TOP RECOMMENDATION")
    L.append(f"  {s['recommendation']}")
    L.append("")
    L.append("Note: all figures are fabricated for demonstration; not from any "
             "employer, customer, or production system.")
    return "\n".join(L)


# ----------------------------------------------------------------- pdf
def build_pdf(s: dict) -> bytes:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    TableStyle)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=LETTER, topMargin=0.7 * inch,
                            bottomMargin=0.7 * inch, leftMargin=0.8 * inch,
                            rightMargin=0.8 * inch)
    styles = getSampleStyleSheet()
    accent = colors.HexColor("#4f5bd5")
    ink = colors.HexColor("#1a1f2e")
    muted = colors.HexColor("#6f7689")

    h = ParagraphStyle("h", parent=styles["Title"], textColor=ink, fontSize=19,
                       spaceAfter=2, alignment=0)
    sub = ParagraphStyle("sub", parent=styles["Normal"], textColor=muted, fontSize=9,
                         spaceAfter=10)
    sec = ParagraphStyle("sec", parent=styles["Heading2"], textColor=accent,
                         fontSize=12, spaceBefore=12, spaceAfter=4)
    body = ParagraphStyle("body", parent=styles["Normal"], textColor=ink, fontSize=10,
                          leading=14)
    status_col = {"Healthy": colors.HexColor("#1f9d57"),
                  "Needs attention": colors.HexColor("#dd9421"),
                  "At risk": colors.HexColor("#e8654f")}.get(s["status"], accent)
    status_style = ParagraphStyle("st", parent=body, textColor=status_col,
                                  fontSize=11, spaceAfter=6)

    el = []
    el.append(Paragraph("Conversational AI — Model Health", h))
    el.append(Paragraph(f"Executive summary &nbsp;|&nbsp; last {s['window_days']} days "
                        f"&nbsp;|&nbsp; {date.today().isoformat()}", sub))
    el.append(Paragraph(f"<b>Overall status: {s['status']}</b>", status_style))
    el.append(Paragraph(s["headline"], body))
    el.append(Spacer(1, 8))

    el.append(Paragraph("Key numbers (latest day)", sec))
    rows = [
        ["Conversations", f"{s['volume']:,}", "Intent accuracy", f"{s['intent_accuracy']*100:.1f}%"],
        ["Resolution rate", f"{s['resolution_rate']*100:.0f}%", "Hallucination risk", f"{s['hallucination_rate']*100:.1f}%"],
        ["Groundedness", f"{s['groundedness_rate']*100:.0f}%", "p95 latency", f"{s['p95_latency_ms']:,} ms"],
        ["Cost / resolved", f"${s['cost_per_resolved']:.4f}", "Intent drift", f"{s['intent_drift']*100:.0f}%"],
    ]
    t = Table(rows, colWidths=[1.5 * inch, 1.4 * inch, 1.5 * inch, 1.4 * inch])
    t.setStyle(TableStyle([
        ("FONTSIZE", (0, 0), (-1, -1), 9.5),
        ("TEXTCOLOR", (0, 0), (0, -1), muted),
        ("TEXTCOLOR", (2, 0), (2, -1), muted),
        ("TEXTCOLOR", (1, 0), (1, -1), ink),
        ("TEXTCOLOR", (3, 0), (3, -1), ink),
        ("FONTNAME", (1, 0), (1, -1), "Helvetica-Bold"),
        ("FONTNAME", (3, 0), (3, -1), "Helvetica-Bold"),
        ("ROWBACKGROUNDS", (0, 0), (-1, -1), [colors.white, colors.HexColor("#f4f4fb")]),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
        ("TOPPADDING", (0, 0), (-1, -1), 7),
        ("LINEBELOW", (0, 0), (-1, -1), 0.4, colors.HexColor("#e7e7f1")),
    ]))
    el.append(t)

    el.append(Paragraph("Health by bucket", sec))
    for b in s["buckets"]:
        el.append(Paragraph(f"<b>{b['name']}.</b> {b['verdict']}", body))
        el.append(Spacer(1, 3))

    el.append(Paragraph("What changed", sec))
    for c in s["changes"]:
        el.append(Paragraph(f"• {c}", body))

    el.append(Paragraph("Top recommendation", sec))
    el.append(Paragraph(s["recommendation"], body))

    el.append(Spacer(1, 14))
    el.append(Paragraph("All figures are fabricated for demonstration; not derived from "
                        "any employer, customer, or production system.",
                        ParagraphStyle("fn", parent=body, fontSize=8, textColor=muted)))

    doc.build(el)
    return buf.getvalue()


# ----------------------------------------------------------------- pptx
def build_pptx(s: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    ACCENT = RGBColor(0x4F, 0x5B, 0xD5)
    INK = RGBColor(0x1A, 0x1F, 0x2E)
    MUTED = RGBColor(0x6F, 0x76, 0x89)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_textbox(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"
        return tb

    def band(slide, color, h=1.1):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        return shp

    # ---- title slide
    s1 = prs.slides.add_slide(blank)
    band(s1, ACCENT, h=7.5)
    add_textbox(s1, 0.9, 2.4, 11.5, 1.2, "Conversational AI — Model Health", 40, WHITE, bold=True)
    add_textbox(s1, 0.9, 3.6, 11.5, 0.6, f"Executive summary  ·  last {s['window_days']} days  ·  {date.today().isoformat()}", 16, WHITE)
    add_textbox(s1, 0.9, 4.4, 11.5, 0.8, f"Overall status: {s['status']}", 20, WHITE, bold=True)
    add_textbox(s1, 0.9, 6.7, 11.5, 0.5, "Fabricated demo data — not from any employer, customer, or production system.", 10, WHITE)

    # ---- summary slide
    s2 = prs.slides.add_slide(blank)
    band(s2, ACCENT)
    add_textbox(s2, 0.7, 0.28, 12, 0.7, "Executive summary", 26, WHITE, bold=True)
    add_textbox(s2, 0.7, 1.4, 11.9, 1.3, s["headline"], 15, INK)
    # key numbers as two columns
    nums = [
        ("Conversations", f"{s['volume']:,}"),
        ("Intent accuracy", f"{s['intent_accuracy']*100:.1f}%"),
        ("Resolution rate", f"{s['resolution_rate']*100:.0f}%"),
        ("Hallucination risk", f"{s['hallucination_rate']*100:.1f}%"),
        ("Groundedness", f"{s['groundedness_rate']*100:.0f}%"),
        ("p95 latency", f"{s['p95_latency_ms']:,} ms"),
        ("Cost / resolved", f"${s['cost_per_resolved']:.4f}"),
        ("Intent drift", f"{s['intent_drift']*100:.0f}%"),
    ]
    y = 2.9
    for i, (label, val) in enumerate(nums):
        col = i % 2
        row = i // 2
        lx = 0.8 + col * 6.2
        ty = y + row * 0.9
        add_textbox(s2, lx, ty, 3.2, 0.5, label, 12, MUTED)
        add_textbox(s2, lx + 3.0, ty, 2.6, 0.5, val, 18, INK, bold=True)

    # ---- one slide per bucket
    for b in s["buckets"]:
        sl = prs.slides.add_slide(blank)
        band(sl, ACCENT)
        add_textbox(sl, 0.7, 0.28, 12, 0.7, b["name"], 26, WHITE, bold=True)
        add_textbox(sl, 0.7, 1.5, 11.9, 1.0, b["verdict"], 16, INK, bold=True)
        ty = 2.9
        for metric, val in b["metrics"]:
            add_textbox(sl, 0.9, ty, 4.5, 0.5, metric, 14, MUTED)
            add_textbox(sl, 5.4, ty, 3, 0.5, val, 16, INK, bold=True)
            ty += 0.62

    # ---- recommendation slide
    sr = prs.slides.add_slide(blank)
    band(sr, ACCENT)
    add_textbox(sr, 0.7, 0.28, 12, 0.7, "What changed & recommendation", 26, WHITE, bold=True)
    ty = 1.6
    add_textbox(sr, 0.8, ty, 11.8, 0.5, "What changed", 16, ACCENT, bold=True)
    ty += 0.7
    for c in s["changes"]:
        add_textbox(sr, 1.0, ty, 11.5, 0.6, "•  " + c, 13, INK)
        ty += 0.62
    ty += 0.3
    add_textbox(sr, 0.8, ty, 11.8, 0.5, "Top recommendation", 16, ACCENT, bold=True)
    ty += 0.7
    add_textbox(sr, 1.0, ty, 11.5, 1.2, s["recommendation"], 15, INK, bold=True)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
